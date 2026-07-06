#!/usr/bin/env python3
"""Phase 1: Convert all source files to Markdown for 期末一键复习（Final Review）."""

import os
import sys
import re
from pathlib import Path

# Fix GBK encoding for Windows console
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

try:
    import docx
except ImportError:
    print("python-docx not installed")
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed")
    sys.exit(1)

try:
    from PyPDF2 import PdfReader
except ImportError:
    print("PyPDF2 not installed")
    sys.exit(1)

BASE_DIR = Path(r"D:\环资法")
OUTPUT_DIR = BASE_DIR / ".temp_work" / "converted"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

EXCLUDE_PATTERNS = [
    "环境法复习资料（完整版）",
    "环境法复习资料（完整版_新）",
    ".temp_work",
]


def should_exclude(filepath: Path) -> bool:
    """Check if a file should be excluded."""
    name = filepath.name
    for pat in EXCLUDE_PATTERNS:
        if pat in name:
            return True
    return False


def convert_docx(filepath: Path) -> str:
    """Convert a .docx file to markdown text."""
    doc = docx.Document(str(filepath))
    lines = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            lines.append("")
            continue
        style = p.style.name
        if style.startswith("Heading"):
            level = style.replace("Heading ", "")
            try:
                level = int(level)
            except ValueError:
                level = 2
            if level > 6:
                level = 2
            lines.append(f"{'#' * level} {text}")
        else:
            lines.append(text)

    # Tables
    for i, table in enumerate(doc.tables):
        lines.append(f"\n**表格 {i+1}**\n")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


def convert_pptx(filepath: Path) -> str:
    """Convert a .pptx/.ppt file to markdown text."""
    prs = Presentation(str(filepath))
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## 幻灯片 {slide_num}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return "\n".join(lines)


def convert_pdf(filepath: Path) -> str:
    """Convert a text-based PDF to markdown text."""
    try:
        reader = PdfReader(str(filepath))
        lines = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text.strip():
                lines.append(f"\n## 第{i}页\n")
                lines.append(text)
        return "\n".join(lines)
    except Exception as e:
        return f"[PDF读取失败: {e}]"


def process_file(filepath: Path, rel_path: str):
    """Process a single file and save as markdown."""
    ext = filepath.suffix.lower()

    if ext == ".docx":
        md_content = convert_docx(filepath)
    elif ext in (".ppt", ".pptx"):
        try:
            md_content = convert_pptx(filepath)
        except Exception as e:
            md_content = f"[PPTX读取失败: {e}]"
    elif ext == ".pdf":
        md_content = convert_pdf(filepath)
    else:
        return False

    # Save output
    safe_name = re.sub(r'[^\w一-鿿\-_]', '_', rel_path)
    if len(safe_name) > 200:
        safe_name = safe_name[:200]
    out_path = OUTPUT_DIR / f"{safe_name}.md"
    out_path.write_text(md_content, encoding="utf-8")

    lines = len(md_content.splitlines())
    chars = len(md_content)
    print(f"  [OK] {rel_path} -> {lines}行, {chars}字符")
    return True


def main():
    total = 0
    success = 0
    errors = []

    # File type priority: notes first, then PPTs, then PDFs
    patterns = [
        ("笔记/*.docx", "笔记"),
        ("笔记/*.pdf", "笔记"),
        ("环境法ppt/*.ppt", "PPT"),
        ("环境法ppt/*.pptx", "PPT"),
        ("往年题/*.docx", "往年题"),
        ("往年题/*.pdf", "往年题"),
        ("*.pdf", "根目录"),
        ("*.docx", "根目录"),
    ]

    for glob_pattern, category in patterns:
        for filepath in sorted(BASE_DIR.glob(glob_pattern)):
            if should_exclude(filepath):
                print(f"  [SKIP] 排除: {filepath.name}")
                continue
            if filepath.stat().st_size < 100:
                print(f"  ⏭️ 跳过小文件: {filepath.name}")
                continue

            rel_path = str(filepath.relative_to(BASE_DIR))
            total += 1
            try:
                result = process_file(filepath, rel_path)
                if result:
                    success += 1
            except Exception as e:
                errors.append(f"{rel_path}: {e}")
                print(f"  [ERR] {rel_path}: {e}")

    print(f"\n{'='*50}")
    print(f"处理完成: {success}/{total} 文件成功转换")
    if errors:
        print(f"错误: {len(errors)}")
        for e in errors:
            print(f"  - {e}")


if __name__ == "__main__":
    main()
