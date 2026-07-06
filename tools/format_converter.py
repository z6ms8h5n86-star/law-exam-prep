#!/usr/bin/env python3
"""
Multi-format → unified Markdown converter for 期末一键复习（Final Review）.

Usage:
    python format_converter.py --input file.docx --mode docx2md
    python format_converter.py --input file.pptx --mode pptx2md
    python format_converter.py --input file.pdf --mode pdf2md
    python format_converter.py --input file.epub --mode epub2md
    python format_converter.py --input file.html --mode html2md
    python format_converter.py --input mubu.json --mode mubu2md
    python format_converter.py --input file.txt --mode txt2md

Output: Markdown text to stdout, or --output to file.
"""

import argparse
import json
import os
import sys
import re
from pathlib import Path

def detect_encoding(filepath: str) -> str:
    """Auto-detect file encoding (UTF-8, GBK, GB2312 fallback)."""
    try:
        import chardet
        with open(filepath, "rb") as f:
            raw = f.read(50000)  # sample first 50KB
        result = chardet.detect(raw)
        return result.get("encoding", "utf-8") or "utf-8"
    except ImportError:
        pass

    # Fallback: try common encodings
    for enc in ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]:
        try:
            with open(filepath, "r", encoding=enc) as f:
                f.read(1000)
            return enc
        except (UnicodeDecodeError, UnicodeError):
            continue
    return "utf-8"

def read_text_file(filepath: str) -> str:
    """Read a text file with auto-detected encoding."""
    enc = detect_encoding(filepath)
    with open(filepath, "r", encoding=enc, errors="replace") as f:
        return f.read()

def docx_to_md(filepath: str) -> str:
    """Convert .docx to Markdown."""
    try:
        from docx import Document
    except ImportError:
        return f"[ERROR: python-docx not installed. Run: pip install python-docx]\n\nFile: {filepath}"

    doc = Document(filepath)
    lines = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        style = para.style.name if para.style else ""

        # Map Word styles to Markdown headings
        if style.startswith("Heading 1") or style == "Title":
            lines.append(f"# {text}")
        elif style.startswith("Heading 2"):
            lines.append(f"## {text}")
        elif style.startswith("Heading 3"):
            lines.append(f"### {text}")
        elif style.startswith("Heading 4"):
            lines.append(f"#### {text}")
        else:
            # Check for bold/italic runs
            formatted = []
            for run in para.runs:
                t = run.text
                if run.bold and run.italic:
                    t = f"***{t}***"
                elif run.bold:
                    t = f"**{t}**"
                elif run.italic:
                    t = f"*{t}*"
                formatted.append(t)
            lines.append("".join(formatted) if formatted else text)

    # Extract tables
    for table in doc.tables:
        lines.append("")
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("|" + "|".join(["---"] * len(cells)) + "|")

    return "\n".join(lines)

def pptx_to_md(filepath: str) -> str:
    """Convert .pptx to Markdown (one slide = one section)."""
    try:
        from pptx import Presentation
    except ImportError:
        return f"[ERROR: python-pptx not installed. Run: pip install python-pptx]\n\nFile: {filepath}"

    prs = Presentation(filepath)
    lines = []
    lines.append(f"# {Path(filepath).stem}\n")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## 第{i}页\n")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level if para.level else 0
                    prefix = "  " * level + "- " if level > 0 else ""
                    # Detect heading-like text
                    if para.runs and para.runs[0].font.size:
                        size_pt = para.runs[0].font.size.pt
                        if size_pt and size_pt >= 28:
                            lines.append(f"### {text}")
                        else:
                            lines.append(f"{prefix}{text}")
                    else:
                        lines.append(f"{prefix}{text}")

            if shape.has_table:
                table = shape.table
                lines.append("")
                for ri, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                    if ri == 0:
                        lines.append("|" + "|".join(["---"] * len(cells)) + "|")
                lines.append("")

        lines.append("")  # blank line between slides

    return "\n".join(lines)

def pdf_to_md(filepath: str) -> str:
    """Convert text-based PDF to Markdown."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return f"[ERROR: PyPDF2 not installed. Run: pip install PyPDF2]\n\nFile: {filepath}"

    try:
        reader = PdfReader(filepath)
    except Exception as e:
        return f"[ERROR: Cannot read PDF: {e}]\n\nFile: {filepath}"

    lines = []
    lines.append(f"# {Path(filepath).stem}\n")

    total_text = ""
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            total_text += text + "\n"

    # If very little text extracted, it might be a scanned PDF
    if len(total_text.strip()) < 100 and len(reader.pages) > 1:
        lines.append("[SCANNED_PDF: This appears to be a scanned/image-based PDF. Use ocr_tool.py for OCR processing.]\n")
        lines.append(f"Pages: {len(reader.pages)}")
        lines.append(f"File: {filepath}")
        return "\n".join(lines)

    # Clean and format extracted text
    # Remove excessive newlines but preserve paragraph breaks
    text = re.sub(r'\n{3,}', '\n\n', total_text)
    text = re.sub(r'(\w)-\n(\w)', r'\1\2', text)  # fix hyphenated line breaks

    lines.append(text)
    return "\n".join(lines)

def epub_to_md(filepath: str) -> str:
    """Convert .epub to Markdown."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        return f"[ERROR: ebooklib or beautifulsoup4 not installed. Run: pip install ebooklib beautifulsoup4]\n\nFile: {filepath}"

    try:
        book = epub.read_epub(filepath)
    except Exception as e:
        return f"[ERROR: Cannot read EPUB: {e}]\n\nFile: {filepath}"

    lines = []
    title = book.get_metadata("DC", "title")
    if title:
        lines.append(f"# {title[0][0]}\n")

    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            # Remove script/style tags
            for tag in soup(["script", "style"]):
                tag.decompose()

            # Convert headings
            for h in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                level = int(h.name[1])
                h.string = f"{'#' * level} {h.get_text().strip()}"

            # Get text with basic formatting
            text = soup.get_text("\n", strip=True)
            # Clean up excessive whitespace
            text = re.sub(r'\n{3,}', '\n\n', text)
            lines.append(text)

    return "\n".join(lines)

def html_to_md(filepath: str) -> str:
    """Convert .html to Markdown."""
    try:
        from markdownify import markdownify as mdify
    except ImportError:
        # Fallback: basic BeautifulSoup extraction
        try:
            from bs4 import BeautifulSoup
            with open(filepath, "r", encoding=detect_encoding(filepath), errors="replace") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()
            return f"# {Path(filepath).stem}\n\n" + soup.get_text("\n", strip=True)
        except ImportError:
            return f"[ERROR: markdownify or beautifulsoup4 not installed. Run: pip install markdownify]\n\nFile: {filepath}"

    enc = detect_encoding(filepath)
    with open(filepath, "r", encoding=enc, errors="replace") as f:
        html_content = f.read()

    md = mdify(html_content, heading_style="ATX", strip=["script", "style", "nav"])
    return f"# {Path(filepath).stem}\n\n{md}"

def mubu_json_to_md(filepath: str) -> str:
    """Convert Mubu (幕布) JSON export to hierarchical Markdown."""
    with open(filepath, "r", encoding="utf-8") as f:
        data = json.load(f)

    lines = []
    title = Path(filepath).stem.replace("_", " ").replace("-", " ")
    lines.append(f"# {title}\n")

    def process_node(node, depth=0):
        """Recursively process Mubu node tree."""
        if isinstance(node, dict):
            text = node.get("text", node.get("title", node.get("name", "")))
            if text:
                prefix = "  " * depth + "- "
                lines.append(f"{prefix}{text.strip()}")

            children = node.get("children", node.get("child", node.get("nodes", [])))
            if isinstance(children, list):
                for child in children:
                    process_node(child, depth + 1)
        elif isinstance(node, list):
            for item in node:
                process_node(item, depth)
        elif isinstance(node, str):
            lines.append(f"{'  ' * depth}- {node.strip()}")

    # Handle different Mubu export structures
    if isinstance(data, dict):
        root_nodes = data.get("data", data.get("nodes", data.get("children", [])))
        if isinstance(root_nodes, list):
            for node in root_nodes:
                process_node(node, 0)
        else:
            process_node(data, 0)
    elif isinstance(data, list):
        for node in data:
            process_node(node, 0)

    return "\n".join(lines)

def txt_to_md(filepath: str) -> str:
    """Convert plain text to basic Markdown."""
    text = read_text_file(filepath)
    title = Path(filepath).stem
    return f"# {title}\n\n{text}"

def main():
    parser = argparse.ArgumentParser(description="Convert various formats to unified Markdown")
    parser.add_argument("--input", "-i", required=True, help="Input file path")
    parser.add_argument("--mode", "-m", choices=[
        "docx2md", "pptx2md", "pdf2md", "epub2md",
        "html2md", "mubu2md", "txt2md", "auto"
    ], default="auto", help="Conversion mode (default: auto-detect by extension)")
    parser.add_argument("--output", "-o", help="Output file path (default: stdout)")
    args = parser.parse_args()

    filepath = args.input
    if not os.path.isfile(filepath):
        print(f"ERROR: File not found: {filepath}", file=sys.stderr)
        sys.exit(1)

    # Auto-detect mode from extension
    mode = args.mode
    if mode == "auto":
        ext = Path(filepath).suffix.lower()
        mode_map = {
            ".docx": "docx2md",
            ".doc": "docx2md",
            ".pptx": "pptx2md",
            ".ppt": "pptx2md",
            ".pdf": "pdf2md",
            ".epub": "epub2md",
            ".html": "html2md",
            ".htm": "html2md",
            ".txt": "txt2md",
            ".md": "txt2md",  # pass-through
            ".json": "mubu2md",  # assume Mubu JSON
        }
        mode = mode_map.get(ext, "txt2md")

    converters = {
        "docx2md": docx_to_md,
        "pptx2md": pptx_to_md,
        "pdf2md": pdf_to_md,
        "epub2md": epub_to_md,
        "html2md": html_to_md,
        "mubu2md": mubu_json_to_md,
        "txt2md": txt_to_md,
    }

    converter = converters.get(mode, txt_to_md)
    try:
        result = converter(filepath)
    except Exception as e:
        result = f"[ERROR during conversion: {e}]\n\nFile: {filepath}"

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(result)
        print(f"Converted: {filepath} → {args.output}", file=sys.stderr)
    else:
        print(result)

if __name__ == "__main__":
    main()
